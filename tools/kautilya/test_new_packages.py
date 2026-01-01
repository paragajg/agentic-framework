"""
Test the newly installed packages: yfinance, python-pptx, python-docx
"""

import sys
from pathlib import Path

try:
    from rich.console import Console
    console = Console()
except ImportError:
    class Console:
        def print(self, *args, **kwargs):
            print(*args)
    console = Console()


def test_yfinance():
    """Test yfinance package."""
    console.print("\n[bold cyan]Test 1: yfinance (Yahoo Finance Data)[/bold cyan]")
    console.print("-" * 60)

    try:
        import yfinance as yf

        # Get stock info (simple test, no network call)
        ticker = yf.Ticker("AAPL")
        console.print(f"✓ yfinance imported successfully")
        console.print(f"  Created ticker object for AAPL")
        console.print(f"  yfinance version: {yf.__version__}")

        return True
    except Exception as e:
        console.print(f"✗ yfinance test failed: {e}")
        return False


def test_python_pptx():
    """Test python-pptx package."""
    console.print("\n[bold cyan]Test 2: python-pptx (PowerPoint)[/bold cyan]")
    console.print("-" * 60)

    try:
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile

        # Create a presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Kautilya AI Framework"
        subtitle.text = "Automated Presentation Generation Test"

        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            temp_file = f.name
            prs.save(temp_file)

        # Verify file exists
        file_size = Path(temp_file).stat().st_size

        console.print(f"✓ python-pptx working correctly")
        console.print(f"  Created PowerPoint presentation")
        console.print(f"  File size: {file_size:,} bytes")
        console.print(f"  Saved to: {temp_file}")

        # Cleanup
        Path(temp_file).unlink()

        return True
    except Exception as e:
        console.print(f"✗ python-pptx test failed: {e}")
        import traceback
        console.print(traceback.format_exc())
        return False


def test_python_docx():
    """Test python-docx package."""
    console.print("\n[bold cyan]Test 3: python-docx (Word Documents)[/bold cyan]")
    console.print("-" * 60)

    try:
        from docx import Document
        import tempfile

        # Create a document
        doc = Document()
        doc.add_heading('Kautilya AI Framework', 0)
        doc.add_paragraph('This is a test document created by the Kautilya AI Framework.')
        doc.add_paragraph('Features:', style='List Bullet')
        doc.add_paragraph('Safe package installation', style='List Bullet')
        doc.add_paragraph('Automated document generation', style='List Bullet')
        doc.add_paragraph('Enterprise-ready security', style='List Bullet')

        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            temp_file = f.name
            doc.save(temp_file)

        # Verify file exists
        file_size = Path(temp_file).stat().st_size

        console.print(f"✓ python-docx working correctly")
        console.print(f"  Created Word document")
        console.print(f"  File size: {file_size:,} bytes")
        console.print(f"  Saved to: {temp_file}")

        # Cleanup
        Path(temp_file).unlink()

        return True
    except Exception as e:
        console.print(f"✗ python-docx test failed: {e}")
        import traceback
        console.print(traceback.format_exc())
        return False


def test_safe_package_manager():
    """Test that packages are whitelisted."""
    console.print("\n[bold cyan]Test 4: SafePackageManager Whitelist[/bold cyan]")
    console.print("-" * 60)

    try:
        from kautilya.safe_package_manager import SafePackageManager

        manager = SafePackageManager()

        packages = ['yfinance', 'python-pptx', 'python-docx']
        all_allowed = True

        for pkg in packages:
            is_allowed, reason = manager.is_package_allowed(pkg)
            status = "✓" if is_allowed else "✗"
            console.print(f"  {status} {pkg}: {reason}")
            if not is_allowed:
                all_allowed = False

        if all_allowed:
            console.print("\n✓ All packages are whitelisted")
            return True
        else:
            console.print("\n✗ Some packages are not whitelisted")
            return False

    except Exception as e:
        console.print(f"✗ SafePackageManager test failed: {e}")
        return False


def main():
    """Run all tests."""
    console.print("\n" + "=" * 60)
    console.print("[bold yellow]New Packages Test Suite[/bold yellow]")
    console.print("=" * 60)

    results = []
    results.append(("yfinance", test_yfinance()))
    results.append(("python-pptx", test_python_pptx()))
    results.append(("python-docx", test_python_docx()))
    results.append(("SafePackageManager", test_safe_package_manager()))

    console.print("\n" + "=" * 60)
    console.print("[bold]Test Results Summary[/bold]")
    console.print("=" * 60)

    for name, result in results:
        status = "[green]✓ PASS[/green]" if result else "[red]✗ FAIL[/red]"
        console.print(f"{status} - {name}")

    all_passed = all(result for _, result in results)

    console.print("\n" + "=" * 60)
    if all_passed:
        console.print("[bold green]✓ ALL TESTS PASSED[/bold green]")
        console.print("=" * 60)
        console.print("\n[cyan]All new packages are installed and working![/cyan]")
        console.print("\nInstalled packages:")
        console.print("  • yfinance (Yahoo Finance data)")
        console.print("  • python-pptx (PowerPoint generation)")
        console.print("  • python-docx (Word document generation)")
        console.print("\nAll packages are whitelisted in SafePackageManager.")
        console.print("Your team can now use them safely!\n")
    else:
        console.print("[bold red]✗ SOME TESTS FAILED[/bold red]")
        console.print("=" * 60 + "\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
